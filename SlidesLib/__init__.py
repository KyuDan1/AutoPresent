from .search import GoogleSearch
from .vqa import VQA
from .image_gen import Dalle3
from .llm import LLM
from .ppt_gen import SlideAgent 
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from mysearchlib import LLM
from pptx.dml.color import RGBColor
def search_result(question: str, screenshot_path: str = "screenshot.png") -> str:
    """
    Search a question on Google, and take a screenshot of the search result.
    Save the screenshot to screenshot_path, and return the path.
    """
    return GoogleSearch.search_result(question, screenshot_path)

def search_image(query: str, save_path: str = 'top_images') -> str:
    """
    Search for an image on Google and download the result to download_path.
    Return download_path.
    """
    return GoogleSearch.search_image(query, save_path)

def get_answer(question: str) -> str:
    """
    Calls the LLM by inputing a question, 
    then get the response of the LLM as the answer
    """
    return LLM.get_answer(question)

def get_code(request:str, examples:str = "") -> str:
    """ 
    Calls the LLM to generate code for a request. 
    request: the task that the model should conduct
    examples: few-shot code examples for the request
    """
    return LLM.get_answer(request, examples)

def generate_image(query: str, save_path: str = "downloaded_image.png") -> str:
    """
    Generate an image based on a text query, save the image to the save_path
    Return the path of the saved image.
    """
    return Dalle3.generate_image(query, save_path)


def add_title(
    slide, text: str, font_size: int = 44, 
    font_color: tuple[int, int, int] = (0, 0, 0),
    background_color: tuple[int, int, int] = None,
):
    """Add a title text to the slide with custom font size and font color (RGB tuple).
    Args:
        slide: Slide object as in pptx library
        text: str, Title text to be added
        font_size: int, Font size in int (point size), e.g., 44
        font_color: tuple(int,int,int), RGB color, e.g., (0, 0, 0)
        background_color: Optional, tuple(int,int,int), RGB color, e.g., (255, 255, 255)
    Rets:
        slide: Slide object with the title added
    """
    title_shape = slide.shapes.title
    if title_shape is None:
        # Add a new text box as the title if no placeholder is found
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(*font_color)
    if background_color is not None:
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = RGBColor(*background_color)
    return slide


def add_text(
    slide, text: str, coords: list[float],
    font_size: int = 20, bold: bool = False,
    color: tuple[int, int, int] = (0, 0, 0),
    background_color: tuple[int, int, int] = None,
    auto_size: bool = True,
):
    """Add a text box at a specified location with custom text and color settings.
    Args:
        slide: Slide object as in pptx library
        text: str, Text to be added
        coords: list(float), [left, top, width, height] in inches
        font_size: int, Font size in int (point size), e.g., 20
        bold: bool, True if bold-type the text, False otherwise
        color: tuple(int,int,int), RGB color, e.g., (0, 0, 0)
        background_color: Optional, tuple(int,int,int), RGB color, e.g., (255, 255, 255)
        auto_size: bool, True if auto-size the text box, False otherwise
    Rets:
        slide: Slide object with the text box added
    """
    # Create the text box shape
    left, top, width, height = coords
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    # Set background color if provided
    if background_color:
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = RGBColor(*background_color)
    else:
        text_box.fill.background()  # No fill if no color is specified

    # Handle line breaks and adjust height
    lines = text.split("\n")
    adjusted_height = height * len(lines)  # Adjust height based on the number of lines
    text_box.height = Inches(adjusted_height)

    # Set text and format it
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    if auto_size:
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Automatically fit the text box to the text

    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return slide


def add_bullet_points(
    slide, bullet_points: list[str], coords: list[float],
    font_size: int = 18, color: tuple[int, int, int] = (0, 0, 0),
    background_color: tuple[int, int, int] = None,
):
    """Add a text box with bullet points.
    Args:
        slide: Slide object as in pptx library
        bullet_points: list(str), List of texts to be added as bullet points
        coords: list(float), [left, top, width, height] in inches
        font_size: int, Font size in int (point size), e.g., 18
        color: tuple(int,int,int), RGB color, e.g., (0, 0, 0)
        background_color: Optional, tuple(int,int,int), RGB color, e.g., (255, 255, 255)
    Rets:
        slide: Slide object with the bullet points added
    """
    left, top, width, height = coords
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    # Set background color if provided
    if background_color:
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = RGBColor(*background_color)
    else:
        text_box.fill.background()  # No fill if no color is specified

    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*color)
        # p.level = bullet_points.index(point)

    return slide


def add_image(slide, image_path: str, coords: list[float]):
    """Add an image in the provided path to the specified coords and sizes.
    Args:
        slide: Slide object as in pptx library
        image_path: str, Path to the image file
        coords: list(float), [left, top, width, height] in inches
    Rets:
        slide: Slide object with the image added
    """
    left, top, width, height = coords
    slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    return slide


def set_background_color(slide, color: tuple[int, int, int]):
    """Set background color for the current slide.
    Args:
        slide: Slide object as in pptx library
        color: tuple(int, int, int), RGB color, e.g., (255, 255, 255)
    Returns:
        modified slide object
    """
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)  # Convert tuple to RGBColor
    return slide